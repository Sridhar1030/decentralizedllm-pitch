#!/usr/bin/env python3
"""Builds the DecentralizedLLM pitch deck (5 slides, 16:9) with python-pptx.

Everything is drawn as native PowerPoint shapes, not images, so the deck stays
editable and stays sharp on a projector. Content lives in CONTENT at the bottom;
the functions above are pure layout.

    pip install python-pptx && python3 build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
GROUND = C(0x0D, 0x14, 0x20)
SURF   = C(0x14, 0x1E, 0x2B)
LINE   = C(0x26, 0x32, 0x3E)
INK    = C(0xE9, 0xEF, 0xF6)
DIM    = C(0x8A, 0xA0, 0xB5)
MUT    = C(0x5D, 0x71, 0x86)
NODE   = [C(0x2D, 0xD4, 0xBF), C(0x8F, 0x8F, 0xFB), C(0xF4, 0x72, 0xA6)]
OK     = C(0x4A, 0xDE, 0x80)
BAD    = C(0xFB, 0x71, 0x85)
ACCENT = C(0xF0, 0xA8, 0x3C)

SANS = "Arial"          # present on macOS and Windows; no silent substitution
MONO = "Courier New"

W, H = 13.333, 7.5

# ---------------------------------------------------------------- primitives
def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def text(slide, s, x, y, w, h, size=14, color=INK, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.15, caps=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = _tf(box); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    r = p.add_run(); r.text = s.upper() if caps else s
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, radius=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        I(x), I(y), I(w), I(h))
    if radius:
        shp.adjustments[0] = radius
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:    shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:    shp.line.fill.background()
    shp.shadow.inherit = False
    _tf(shp)
    return shp

def hrule(slide, x, y, w, color=LINE, lw=1.0):
    ln = slide.shapes.add_connector(1, I(x), I(y), I(x + w), I(y))
    ln.line.color.rgb = color; ln.line.width = Pt(lw)
    return ln

def slide_base(prs, kicker, title, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, fill=GROUND)
    text(s, kicker, 0.75, 0.52, 8.0, 0.3, size=10.5, color=ACCENT, font=MONO, caps=True)
    text(s, title, 0.75, 0.85, 11.0, 0.85, size=33, color=INK, bold=True)
    hrule(s, 0.75, 1.82, W - 1.5)
    text(s, f"{n} / 5", W - 1.55, 0.52, 0.8, 0.3, size=10.5, color=MUT,
         font=MONO, align=PP_ALIGN.RIGHT)
    return s

def bullets(slide, items, x, y, w, size=15, gap=0.46, dot=ACCENT):
    """items: list of (bold_lead, rest). bold_lead may be ''."""
    for i, (lead, rest) in enumerate(items):
        yy = y + i * gap
        rect(slide, x, yy + 0.085, 0.075, 0.075, fill=dot)
        box = slide.shapes.add_textbox(I(x + 0.28), I(yy - 0.03), I(w - 0.28), I(gap))
        tf = _tf(box); p = tf.paragraphs[0]; p.line_spacing = 1.2
        if lead:
            r = p.add_run(); r.text = lead
            r.font.size = Pt(size); r.font.bold = True; r.font.name = SANS
            r.font.color.rgb = INK
        if rest:
            r = p.add_run(); r.text = (" " if lead else "") + rest
            r.font.size = Pt(size); r.font.name = SANS; r.font.color.rgb = DIM

def metrics(slide, cards, x, y, w, h=1.25, gapx=0.22):
    """cards: list of (value, label, color)."""
    n = len(cards)
    cw = (w - gapx * (n - 1)) / n
    for i, (val, label, col) in enumerate(cards):
        cx = x + i * (cw + gapx)
        rect(slide, cx, y, cw, h, fill=SURF, line=LINE, radius=0.06)
        text(slide, val, cx + 0.22, y + 0.20, cw - 0.44, 0.5,
             size=25, color=col, font=MONO, bold=True)
        text(slide, label, cx + 0.22, y + 0.78, cw - 0.44, 0.42,
             size=9.5, color=DIM, font=MONO, spacing=1.1)

def table(slide, head, rows, x, y, w, colw, rowh=0.42, size=12):
    """head: list of column titles. rows: list of (label, *cells, colorlist)."""
    cx = x
    for i, htxt in enumerate(head):
        text(slide, htxt, cx, y, colw[i], 0.3, size=9, color=MUT, font=MONO, caps=True,
             align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
        cx += colw[i]
    hrule(slide, x, y + 0.34, w)
    for r, (cells, colors) in enumerate(rows):
        yy = y + 0.46 + r * rowh
        cx = x
        for i, cell in enumerate(cells):
            text(slide, cell, cx, yy, colw[i], rowh, size=size,
                 color=colors[i], font=SANS if i == 0 else MONO,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
            cx += colw[i]
        hrule(slide, x, yy + 0.3, w, color=C(0x1C, 0x27, 0x33))

def pipeline(slide, x, y, w, labels, sublabels, packet=None, caption=None):
    """The three-shard pipeline drawn as native shapes."""
    bw, bh, gap = 2.05, 1.05, 0.52
    coordw = 1.35
    rect(slide, x, y + 0.1, coordw, bh - 0.2, fill=SURF, line=LINE, radius=0.08)
    text(slide, "coordinator", x + 0.12, y + 0.48, coordw - 0.24, 0.3,
         size=9, color=MUT, font=MONO, align=PP_ALIGN.CENTER)
    bx = x + coordw + gap
    for i in range(3):
        cx = bx + i * (bw + gap)
        rect(slide, cx, y, bw, bh, fill=SURF, line=NODE[i], radius=0.08)
        text(slide, labels[i], cx + 0.16, y + 0.16, bw - 0.32, 0.3,
             size=11.5, color=NODE[i], font=MONO, bold=True)
        text(slide, sublabels[i], cx + 0.16, y + 0.52, bw - 0.32, 0.45,
             size=9, color=DIM, font=MONO, spacing=1.15)
        # wire into this node
        hrule(slide, cx - gap, y + bh / 2, gap, color=LINE, lw=1.5)
    endx = bx + 3 * (bw + gap)
    hrule(slide, endx - gap, y + bh / 2, gap, color=LINE, lw=1.5)
    text(slide, "token", endx, y + bh / 2 - 0.14, 1.1, 0.3,
         size=10, color=OK, font=MONO)
    if packet:
        px = bx + bw + gap * 0.18
        rect(slide, px, y + bh / 2 - 0.09, 0.34, 0.18, fill=NODE[0], radius=0.3)
        text(slide, packet, px - 0.45, y + bh / 2 - 0.48, 1.3, 0.3,
             size=9, color=NODE[0], font=MONO, align=PP_ALIGN.CENTER)
    if caption:
        text(slide, caption, x, y + bh + 0.28, w, 0.3, size=10, color=MUT, font=MONO)

def footnote(slide, s):
    text(slide, s, 0.75, H - 0.72, W - 1.5, 0.45, size=9, color=MUT,
         font=MONO, spacing=1.25)

def speaker_notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


# ================================================================ slide bodies
# Content is verbatim from knowledge-base/40-PITCH.md section 3. Every number
# carries its (measured)/(derived)/(modelled) tag exactly as the claims ledger
# (decisions/ADR-013) requires. Do not "tidy" a tag away.

def slide1(prs):
    s = slide_base(prs, "the problem", "Why split a model at all", 1)
    bullets(s, [
        ("Five hospitals, one shared model;", "no member may hold it."),
        ("70B fp16 across 3 nodes:", "47 GB per shard — fits nothing."),
        ("Across 16 nodes: 8.8 GB", "— fits an ordinary laptop."),
        ("Below ~13B this is theatre.", "The memory wall is the product."),
        ("Sovereign cloud IaaS, Europe:", "$12.6B in 2026, from $6.9B."),
    ], 0.75, 2.15, 5.1, size=13.5, gap=0.60)
    text(s, "Gartner, 9 Feb 2026 (cited)", 1.03, 5.20, 4.8, 0.3,
         size=9, color=MUT, font=MONO)

    # --- bar chart: GB per shard, linear axis to 50 GB
    bx, by, bw = 6.45, 2.30, 5.55
    axis = 50.0
    bars = [("70B fp16, N=3",  47.0, False),
            ("70B fp16, N=8",  17.6, False),
            ("70B int4, N=3",  11.8, True),
            ("70B fp16, N=16",  8.8, True)]
    rowh = 0.78
    for i, (lab, gb, fits) in enumerate(bars):
        yy = by + i * rowh
        text(s, lab, bx, yy, 1.85, 0.3, size=10, color=DIM, font=MONO)
        wpx = (gb / axis) * (bw - 1.95)
        rect(s, bx + 1.95, yy + 0.02, wpx, 0.30,
             fill=NODE[0] if fits else C(0x2A, 0x35, 0x42),
             line=None if fits else MUT, radius=0.15)
        text(s, f"{gb:.1f} GB", bx + 1.95 + wpx + 0.10, yy + 0.02, 1.2, 0.3,
             size=10.5, color=OK if fits else MUT, font=MONO, bold=fits)
    # device ceilings, drawn through the bars
    for gb, lab in ((16, "16 GB laptop"), (24, "24 GB 4090")):
        cx = bx + 1.95 + (gb / axis) * (bw - 1.95)
        ln = s.shapes.add_connector(1, I(cx), I(by - 0.22), I(cx), I(by + 4 * rowh - 0.10))
        ln.line.color.rgb = ACCENT; ln.line.width = Pt(1.0); ln.line.dash_style = 4
        text(s, lab, cx - 0.55, by - 0.52, 1.6, 0.3, size=8.5, color=ACCENT,
             font=MONO, align=PP_ALIGN.CENTER)
    text(s, "80 GB H100 — off scale →", bx + bw - 2.1, by - 0.52, 2.1, 0.3,
         size=8.5, color=MUT, font=MONO, align=PP_ALIGN.RIGHT)
    text(s, "derived from Llama-3.3-70B config shapes — 70.552B params",
         bx, by + 4 * rowh + 0.18, bw, 0.3, size=9, color=MUT, font=MONO)

    footnote(s, "Left of the 16 GB line, a room of ordinary laptops holds a frontier-class "
                "model. Right of it, they cannot — and quantization does not fix it: int4 of a 7B is still a 7B.")
    speaker_notes(s,
        "Open with the hospital sentence — a named buyer inside ten seconds, before any "
        "architecture. Then the chart, and be blunt about what it says: at half a billion parameters, "
        "splitting a model is theatre, and our own demo proves it by loading the whole thing on every "
        "node. The interesting line is the 16 GB one. Left of it, a room of ordinary laptops can hold a "
        "frontier-class model. Right of it, they cannot, and no amount of quantization fixes that — "
        "int4 of a 7B is still a 7B. That boundary is the entire product.")
    return s


def slide2(prs):
    s = slide_base(prs, "what runs today", "One token, three hops, today", 2)
    bullets(s, [
        ("Qwen2.5-0.5B, 24 layers,", "split 8/8/8 across three containers."),
        ("FastAPI nodes, CPU-only, 2 vCPU;", "gateway, coordinator, Prom, Grafana."),
        ("Node2 owns lm_head,", "so node2 is where a token is born."),
        ("Pipeline parallelism is not ours:", "GPipe 2019, Petals 2022."),
        ("What is new:", "the shard boundary is a trust boundary."),
    ], 0.75, 2.15, 11.8, size=14, gap=0.50)

    pipeline(s, 0.75, 5.02, 11.8,
             ["node0", "node1", "node2"],
             ["embed + L0–7\n8.00 layer-eq",
              "L8–15\n8.00 layer-eq",
              "L16–23 + norm + lm_head\n17.13 layer-eq"],
             packet="[seq, 896] fp32 · base64 · JSON",
             caption="client → gateway :8080 (api-key, circuit breaker) → coordinator "
                     ":8081 → node0 → node1 → node2 → token")
    # lm_head callout, in --bad, against node2
    rect(s, 9.55, 4.58, 3.0, 0.34, fill=None, line=BAD, radius=0.10)
    text(s, "lm_head = 9.13 layer-equivalents", 9.67, 4.64, 2.8, 0.3,
         size=9.5, color=BAD, font=MONO, bold=True)

    footnote(s, "runs today — `docker compose up`.  The \"equal\" 8/8/8 split really runs "
                "8/8/17: 1.539x off balance on layer-equivalents, 1.30x measured wall clock (90-AUDIT F11).")
    speaker_notes(s,
        "This runs. Three containers, CPU only, no GPUs, and a real completion comes back. Say the "
        "prior art out loud before anyone asks: GPipe defined pipeline parallelism in 2019, Petals "
        "shipped it decentralized in 2022, and vLLM ships it today behind one flag. We did not invent "
        "the mechanism, and claiming we did is how you lose a technical judge in the first minute. What "
        "is different is that in all three of those, the shard boundary is a performance boundary. Here "
        "it is a trust boundary — and that changes what you build.")
    return s


def slide3(prs):
    s = slide_base(prs, "the engineering", "Three bets: wire, bytes, queue", 3)
    bullets(s, [
        ("Wire:", "8.483 → 0.089 ms per hop (reported, T1-A4)."),
        ("4.0 ms/hop of that", "was TLS cert parsing on plain http (reported)."),
        ("Bytes:", "bf16 halves the wire, 99.41% top-1, 3.5 µs (measured)."),
        ("One channel is 972x the median", "— naive int8 outputs garbage (measured)."),
        ("Queue:", "U = min(1, R/S). One request cannot fill three stages (derived)."),
    ], 0.75, 2.10, 11.8, size=13.5, gap=0.44)

    py = 4.50
    # -- Panel A: wire
    text(s, "A · WIRE", 0.75, py, 1.5, 0.3, size=9.5, color=NODE[0], font=MONO, bold=True)
    rect(s, 0.75, py + 0.34, 2.55, 0.26, fill=C(0x2A, 0x35, 0x42), radius=0.12)
    text(s, "HTTP+JSON+b64  8.483 ms", 0.75, py + 0.66, 3.4, 0.3, size=9, color=MUT, font=MONO)
    rect(s, 0.75, py + 1.00, 0.14, 0.26, fill=NODE[0], radius=0.12)
    text(s, "DLP frame  0.089 ms", 0.98, py + 1.02, 3.4, 0.3, size=9, color=NODE[0], font=MONO)
    text(s, "95x", 3.42, py + 0.55, 0.8, 0.4, size=19, color=INK, font=MONO, bold=True)
    text(s, "40-byte header · persistent socket · TCP_NODELAY",
         0.75, py + 1.36, 3.7, 0.3, size=8.5, color=MUT, font=MONO)

    # -- Panel B: bytes staircase
    text(s, "B · BYTES", 4.55, py, 1.5, 0.3, size=9.5, color=NODE[1], font=MONO, bold=True)
    steps = [("fp32", 3584, 0.90), ("bf16", 1792, 0.62), ("int8+outliers", 906, 0.40)]
    for i, (lab, b, hgt) in enumerate(steps):
        sx = 4.55 + i * 1.32
        rect(s, sx, py + 1.36 - hgt, 0.95, hgt,
             fill=NODE[1] if i < 2 else C(0x2A, 0x35, 0x42),
             line=None if i < 2 else MUT, radius=0.05)
        text(s, f"{b} B", sx, py + 1.40, 1.1, 0.3, size=9, color=INK if i < 2 else MUT, font=MONO)
        text(s, lab, sx, py + 1.62, 1.32, 0.3, size=8.5, color=DIM, font=MONO)
    text(s, "✕", 6.95, py + 0.98, 0.5, 0.4, size=18, color=BAD, font=MONO, bold=True)
    text(s, "LAN crossover: ~2 Mbit/s — stop at bf16",
         4.55, py + 0.28, 3.6, 0.3, size=8.5, color=BAD, font=MONO)

    # -- Panel C: queue gantt
    text(s, "C · QUEUE", 8.85, py, 1.5, 0.3, size=9.5, color=NODE[2], font=MONO, bold=True)
    for row, (lab, filled, util) in enumerate((("R=1", False, "33%"), ("R=3", True, "~100%"))):
        gy = py + 0.36 + row * 0.62
        text(s, lab, 8.85, gy + 0.10, 0.5, 0.3, size=9, color=DIM, font=MONO)
        for lane in range(3):
            ly = gy + lane * 0.13
            rect(s, 9.42, ly, 2.55, 0.10, fill=C(0x1C, 0x27, 0x33))
            if filled:
                rect(s, 9.42, ly, 2.55, 0.10, fill=NODE[2])
            else:                       # diagonal staircase: one lane lit at a time
                rect(s, 9.42 + lane * 0.85, ly, 0.85, 0.10, fill=NODE[2])
        text(s, util, 12.05, gy + 0.10, 0.9, 0.3, size=11, color=INK, font=MONO, bold=True)
    text(s, "1.40 → 4.21 tok/s at unchanged per-request latency (modelled)",
         8.85, py + 1.62, 4.2, 0.3, size=8.5, color=MUT, font=MONO)

    footnote(s, "Panel B is a negative result we kept on the slide: below bf16 nothing pays until the "
                "link drops under ~2 Mbit/s (90-AUDIT F01 supersedes the team file here). Panel A is reported "
                "by T1-A4 and has no script in bench/ - reported, not reproduced here.")
    speaker_notes(s,
        "Three bets, and one of them is a negative result we kept on the slide. The wire: ninety-five "
        "times per hop, and the funny part is that four milliseconds of every hop was parsing X.509 "
        "certificates for a connection that never used TLS. The bytes: bf16 is free and we stop there, "
        "because we measured the crossover — below bf16 you only win under about two "
        "megabit. The queue is the one people miss: a single request cannot fill a pipeline, because "
        "token t+1 waits on token t. Only concurrency can, and three is the knee.")
    return s


def slide4(prs):
    s = slide_base(prs, "the scoreboard", "The scoreboard, and the caveats", 4)
    bullets(s, [
        ("1.27 → 24.2 tok/s,", "19x, same three containers."),
        ("147,200 → 543 forwards/node:", "271x redundant compute."),
        ("1,821.7 MB → 1.948 MB:", "935x fewer bytes."),
        ("Caveat: bytes, not seconds.", "On a fast LAN we are compute-bound."),
        ("Caveat:", "v1 has never run as an integrated system."),
    ], 0.75, 2.15, 4.35, size=12.5, gap=0.62)

    rows = [
        (("bytes on the wire",        "1,821.7 MB", "1.95 MB",    "935x",     "derived"),
         (INK, MUT, INK, OK, MUT)),
        (("redundant position-fwds",  "147,200",    "543",        "271x",     "derived"),
         (INK, MUT, INK, OK, MUT)),
        (("return path / token",      "607,744 B",  "4 B",        "151,936x", "derived"),
         (INK, MUT, INK, OK, MUT)),
        (("throughput, R=3 balanced", "1.27 tok/s", "24.21 tok/s","19.0x",    "modelled"),
         (INK, MUT, INK, ACCENT, MUT)),
        (("pipeline utilisation",     "33%",        "~100%",      "3x",       "derived"),
         (INK, MUT, INK, OK, MUT)),
    ]
    table(s, ["", "v0", "v1", "factor", "tag"], rows,
          5.55, 2.15, 7.03, [2.55, 1.42, 1.42, 1.06, 0.58], rowh=0.44, size=11.5)

    # caveat box, SAME font size as the table body (11.5) - deliberate, see 40-PITCH
    rect(s, 0.75, 5.62, 11.83, 0.92, fill=None, line=BAD, lw=1.5, radius=0.05)
    text(s, "935x is wire bytes, not wall clock, and it is conservative: v0 star-routes, so the true figure "
            "is ~1,657x. 19x = 6.3x single-stream x 3.0x concurrency. Never run end-to-end.",
         1.00, 5.80, 11.35, 0.6, size=11.5, color=BAD, spacing=1.25)
    speaker_notes(s,
        "The caveat is in the same size type as the numbers, on purpose. If I shrink it, you stop "
        "believing the table. Two-seventy-one-x is the number I would defend hardest — it is pure "
        "arithmetic on the model config, not a benchmark. Nine-hundred-thirty-five-x is the biggest "
        "number and the weakest claim, because it is bytes and on this LAN we are compute-bound; it wins "
        "on WAN, on one-gig ethernet, and at long context. Nineteen-x is modelled from stage times we "
        "measured on one laptop. We have not run v1 end to end. That is the honest state.")
    return s


def slide5(prs):
    s = slide_base(prs, "the honest close", "The memory wall, and the ask", 5)

    text(s, "IS", 0.75, 2.15, 3.0, 0.4, size=20, color=OK, font=MONO, bold=True)
    hrule(s, 0.75, 2.60, 5.6, color=OK)
    for i, (line, sub) in enumerate((
        ("The option when weights exceed", "the device you are allowed to use."),
        ("70B int4, N=3 → 11.8 GB/shard", "the memory wall, in one number."),
        ("A trust boundary at the shard cut", "not merely a performance boundary."),
    )):
        yy = 2.78 + i * 0.86
        text(s, "✓", 0.75, yy, 0.3, 0.3, size=13, color=OK, font=MONO, bold=True)
        text(s, line, 1.12, yy, 5.2, 0.3, size=13, color=INK)
        text(s, sub, 1.12, yy + 0.28, 5.2, 0.3, size=10.5, color=MUT, font=MONO)

    text(s, "IS NOT", 6.95, 2.15, 3.0, 0.4, size=20, color=BAD, font=MONO, bold=True)
    hrule(s, 6.95, 2.60, 5.63, color=BAD)
    for i, (line, sub) in enumerate((
        ("Cheaper than an API", "$5.76/day electricity vs $1.73/day of tokens (modelled)."),
        ("Encryption", "~35% of tokens recoverable from an 8-layer activation (cited)."),
        ("True yet, in our own repo", "node.py:36 loads the whole checkpoint today."),
    )):
        yy = 2.78 + i * 0.86
        text(s, "✕", 6.95, yy, 0.3, 0.3, size=13, color=BAD, font=MONO, bold=True)
        text(s, line, 7.32, yy, 5.2, 0.3, size=13, color=INK)
        text(s, sub, 7.32, yy + 0.28, 5.3, 0.3, size=10.5, color=MUT, font=MONO)

    # roadmap strip
    hrule(s, 0.75, 5.62, 11.83)
    stops = [("pre-sliced shards", "v1, days"), ("KV cache + local argmax", "v1, days"),
             ("binary bf16 frame", "v1, days"), ("reshard-on-failure", "~63 net LOC")]
    for i, (lab, tag) in enumerate(stops):
        sx = 0.75 + i * 2.28
        rect(s, sx, 5.55, 0.13, 0.13, fill=ACCENT, radius=0.5)
        text(s, lab, sx, 5.76, 2.2, 0.3, size=10.5, color=INK, font=MONO)
        text(s, tag, sx, 5.98, 2.2, 0.3, size=9, color=MUT, font=MONO)

    rect(s, 9.10, 6.32, 3.48, 0.78, fill=SURF, line=ACCENT, lw=1.5, radius=0.06)
    text(s, "ASK", 9.32, 6.42, 1.0, 0.25, size=9, color=ACCENT, font=MONO, caps=True)
    text(s, "<FILL: pilot partner / judges' pick / compute credits>",
         9.32, 6.64, 3.1, 0.4, size=11, color=INK, bold=True, spacing=1.1)
    speaker_notes(s,
        "I want to end on what this is not, because every line is a question you were going to ask. It "
        "is not cheaper than an API — a sixteen-node fleet burns five seventy-six a day in "
        "electricity to make a dollar seventy-three of tokens. It is not encryption; split inference "
        "raises the cost of recovering a prompt, it does not prevent it. And today every node still "
        "loads the whole checkpoint, which is forty lines from being false. What it is, is the only "
        "option when the weights do not fit the machine you are allowed to use.")
    return s


def build(path=None):
    # write next to this script, not into whatever cwd it was invoked from
    path = path or os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "DecentralizedLLM-deck.pptx")
    prs = Presentation()
    prs.slide_width, prs.slide_height = I(W), I(H)
    for fn in (slide1, slide2, slide3, slide4, slide5):
        fn(prs)
    prs.save(path)
    return prs, path


if __name__ == "__main__":
    prs, path = build()
    # ponytail: the only check worth having - the deck exists, is 5 slides, and
    # every slide carries speaker notes. Layout is verified by looking at it.
    assert len(prs.slides) == 5, f"expected 5 slides, got {len(prs.slides)}"
    for i, sl in enumerate(prs.slides, 1):
        assert sl.notes_slide.notes_text_frame.text.strip(), f"slide {i} has no speaker notes"
    print(f"ok: {path} — {len(prs.slides)} slides, notes on all")
